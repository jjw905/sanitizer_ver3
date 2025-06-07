# utils/office_reconstructor.py - Office 문서 콘텐츠 재조립

import os
import zipfile
import shutil
from xml.etree import ElementTree as ET
from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from PIL import Image
import io
import config


def reconstruct_office_document(file_path: str, output_dir: str = None) -> tuple[str, dict]:
    """Office 문서에서 안전한 콘텐츠만 추출하여 재조립"""
    if output_dir is None:
        output_dir = config.DIRECTORIES['sanitized_output']

    ext = os.path.splitext(file_path)[1].lower()
    file_name = os.path.splitext(os.path.basename(file_path))[0]
    extracted_content = {'text': '', 'images': [], 'tables': []}

    try:
        if ext in ('.docx', '.docm'):
            clean_file, extracted_content = reconstruct_word_document(file_path, output_dir, file_name)
        elif ext in ('.xlsx', '.xlsm'):
            clean_file, extracted_content = reconstruct_excel_document(file_path, output_dir, file_name)
        elif ext in ('.pptx', '.pptm'):
            clean_file, extracted_content = reconstruct_powerpoint_document(file_path, output_dir, file_name)
        else:
            return None, None

        return clean_file, extracted_content

    except Exception as e:
        print(f"Office 문서 재조립 오류: {e}")
        return None, None


def reconstruct_word_document(file_path: str, output_dir: str, file_name: str) -> tuple[str, dict]:
    """Word 문서 재조립"""
    extracted_content = {'text': '', 'images': [], 'paragraphs': 0}

    try:
        # 원본 문서 열기
        doc = Document(file_path)

        # 새 문서 생성
        new_doc = Document()

        # 텍스트 추출 및 재구성
        all_text = []
        for paragraph in doc.paragraphs:
            text = paragraph.text.strip()
            if text:
                new_doc.add_paragraph(text)
                all_text.append(text)
                extracted_content['paragraphs'] += 1

        extracted_content['text'] = '\n'.join(all_text)

        # 테이블 추출 및 재구성
        for table in doc.tables:
            rows = len(table.rows)
            cols = len(table.columns) if table.rows else 0

            if rows > 0 and cols > 0:
                new_table = new_doc.add_table(rows=rows, cols=cols)
                new_table.style = 'Table Grid'

                for i, row in enumerate(table.rows):
                    for j, cell in enumerate(row.cells):
                        new_table.rows[i].cells[j].text = cell.text.strip()

                extracted_content['tables'].append(f"{rows}x{cols}")

        # 이미지 추출 (안전하게)
        try:
            # 임시 압축 해제
            temp_dir = "temp_word_extract"
            shutil.rmtree(temp_dir, ignore_errors=True)

            with zipfile.ZipFile(file_path, 'r') as zip_ref:
                zip_ref.extractall(temp_dir)

            # media 폴더에서 이미지 찾기
            media_dir = os.path.join(temp_dir, "word", "media")
            if os.path.exists(media_dir):
                for img_file in os.listdir(media_dir):
                    if img_file.lower().endswith(('.png', '.jpg', '.jpeg', '.gif', '.bmp')):
                        img_path = os.path.join(media_dir, img_file)

                        # 이미지 검증
                        try:
                            with Image.open(img_path) as img:
                                img.verify()
                            extracted_content['images'].append(img_file)

                            # 새 문서에 이미지 추가
                            new_doc.add_paragraph()
                            new_doc.add_picture(img_path,
                                                width=new_doc.sections[0].page_width - new_doc.sections[0].left_margin -
                                                      new_doc.sections[0].right_margin)
                        except:
                            pass

            shutil.rmtree(temp_dir, ignore_errors=True)

        except Exception:
            pass

        # 새 문서 저장
        os.makedirs(output_dir, exist_ok=True)
        clean_file = os.path.join(output_dir, f"{file_name}_reconstructed.docx")
        new_doc.save(clean_file)

        return clean_file, extracted_content

    except Exception as e:
        print(f"Word 재조립 오류: {e}")
        return None, None


def reconstruct_excel_document(file_path: str, output_dir: str, file_name: str) -> tuple[str, dict]:
    """Excel 문서 재조립"""
    extracted_content = {'text': '', 'sheets': 0, 'cells': 0}

    try:
        # 원본 파일 열기
        from openpyxl import load_workbook
        wb = load_workbook(file_path, data_only=True)

        # 새 워크북 생성
        new_wb = Workbook()
        new_wb.remove(new_wb.active)  # 기본 시트 제거

        all_text = []

        # 각 시트 처리
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            new_sheet = new_wb.create_sheet(title=sheet_name[:31])  # Excel 시트 이름 제한
            extracted_content['sheets'] += 1

            # 셀 데이터 복사 (수식 제외, 값만)
            for row in sheet.iter_rows():
                for cell in row:
                    if cell.value is not None:
                        new_cell = new_sheet.cell(row=cell.row, column=cell.column)
                        new_cell.value = cell.value

                        # 기본 서식만 복사
                        if cell.font:
                            new_cell.font = cell.font.copy()
                        if cell.fill:
                            new_cell.fill = cell.fill.copy()
                        if cell.border:
                            new_cell.border = cell.border.copy()

                        extracted_content['cells'] += 1
                        all_text.append(str(cell.value))

            # 열 너비 조정
            for column in sheet.columns:
                col_letter = column[0].column_letter
                if sheet.column_dimensions[col_letter].width:
                    new_sheet.column_dimensions[col_letter].width = sheet.column_dimensions[col_letter].width

        extracted_content['text'] = ' '.join(all_text[:1000])  # 처음 1000개 셀만

        # 새 파일 저장
        os.makedirs(output_dir, exist_ok=True)
        clean_file = os.path.join(output_dir, f"{file_name}_reconstructed.xlsx")
        new_wb.save(clean_file)

        return clean_file, extracted_content

    except Exception as e:
        print(f"Excel 재조립 오류: {e}")
        return None, None


def reconstruct_powerpoint_document(file_path: str, output_dir: str, file_name: str) -> tuple[str, dict]:
    """PowerPoint 문서 재조립"""
    extracted_content = {'text': '', 'slides': 0, 'images': []}

    try:
        # 원본 프레젠테이션 열기
        prs = Presentation(file_path)

        # 새 프레젠테이션 생성
        new_prs = Presentation()

        # 기본 빈 슬라이드 제거
        if len(new_prs.slides) > 0:
            rId = new_prs.slides._sldIdLst[0].rId
            new_prs.part.drop_rel(rId)
            del new_prs.slides._sldIdLst[0]

        all_text = []

        # 각 슬라이드 처리
        for slide in prs.slides:
            # 빈 슬라이드 레이아웃 사용
            slide_layout = new_prs.slide_layouts[5]  # blank layout
            new_slide = new_prs.slides.add_slide(slide_layout)
            extracted_content['slides'] += 1

            # 텍스트 추출 및 재구성
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    # 텍스트 박스 추가
                    left = shape.left if hasattr(shape, 'left') else 1000000
                    top = shape.top if hasattr(shape, 'top') else 1000000
                    width = shape.width if hasattr(shape, 'width') else 5000000
                    height = shape.height if hasattr(shape, 'height') else 1000000

                    textbox = new_slide.shapes.add_textbox(left, top, width, height)
                    text_frame = textbox.text_frame
                    text_frame.text = shape.text

                    slide_text.append(shape.text)

                # 이미지 처리 (안전하게)
                if hasattr(shape, "image") and shape.shape_type == 13:  # Picture type
                    try:
                        # 이미지 데이터 추출
                        image = shape.image
                        image_bytes = image.blob

                        # 이미지 검증
                        img = Image.open(io.BytesIO(image_bytes))
                        img.verify()

                        # 이미지 재추가
                        left = shape.left if hasattr(shape, 'left') else 1000000
                        top = shape.top if hasattr(shape, 'top') else 1000000

                        pic = new_slide.shapes.add_picture(
                            io.BytesIO(image_bytes),
                            left, top,
                            width=shape.width if hasattr(shape, 'width') else None,
                            height=shape.height if hasattr(shape, 'height') else None
                        )

                        extracted_content['images'].append(f"Slide{extracted_content['slides']}_Image")
                    except:
                        pass

            all_text.extend(slide_text)

        extracted_content['text'] = '\n'.join(all_text)

        # 새 파일 저장
        os.makedirs(output_dir, exist_ok=True)
        clean_file = os.path.join(output_dir, f"{file_name}_reconstructed.pptx")
        new_prs.save(clean_file)

        return clean_file, extracted_content

    except Exception as e:
        print(f"PowerPoint 재조립 오류: {e}")
        return None, None